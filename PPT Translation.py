from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from concurrent.futures import ThreadPoolExecutor, as_completed
import sys
import warnings
from pptx.dml.color import RGBColor
from pptx.util import Pt
import re
from datetime import datetime
import os
import deepl
 
DEEPL_API_KEY = "04a5a7ce-543e-447f-a2b3-1a8f0bdc9e2a:fx"  # <-- Replace with your actual key
deepl_translator = deepl.Translator(DEEPL_API_KEY)
 
warnings.filterwarnings('ignore', message='Could not copy all formatting properties')
 
extracted_texts = []
translated_texts = []
translation_log = []
 
COMPANY_NAME_PATTERNS = [
    r'\bpromptor[a-z]*\.ai\b',
    r'\bpromptor[a-z]*\s*\.\s*ai\b',
    r'\bpromptor[a-z]*\s+ai\b',
]
 
PLACEHOLDER_TOKEN = "XCOMPANYX"
CORRECT_COMPANY_NAME = "Promptora.ai"
MIN_FONT_SIZE_PT = 7.5
 
def split_text_into_chunks(text, max_length=2000):
    return [text[i:i+max_length] for i in range(0, len(text), max_length)]
 
def translate_chunk(chunk, source_lang, target_lang):
    # DeepL uses language codes like 'EN', 'FR', 'DE', etc.
    # source_lang can be None for auto-detect
    try:
        result = deepl_translator.translate_text(
            chunk,
            source_lang=source_lang.upper() if source_lang and source_lang != 'auto' else None,
            target_lang=target_lang.upper()
        )
        return result.text
    except Exception as e:
        print(f"DeepL translation error: {e}", file=sys.stderr)
        return chunk  # fallback to original if error
 
def protect_company_name(text):
    if not text or not text.strip():
        return text
    protected_text = text
    for pattern in COMPANY_NAME_PATTERNS:
        protected_text = re.sub(pattern, PLACEHOLDER_TOKEN, protected_text, flags=re.IGNORECASE)
    return protected_text
 
def restore_company_name(text):
    if not text:
        return text
    restored_text = text.replace(PLACEHOLDER_TOKEN, CORRECT_COMPANY_NAME)
    for pattern in COMPANY_NAME_PATTERNS:
        restored_text = re.sub(pattern, CORRECT_COMPANY_NAME, restored_text, flags=re.IGNORECASE)
    return restored_text
 
def translate_text_safely(text, source_lang, target_lang):
    if not text or not text.strip():
        extracted_texts.append(text)
        translated_texts.append(text)
        translation_log.append("Skipped (empty or whitespace)")
        return text
 
    original_text = text
    extracted_texts.append(original_text)
    protected_text = protect_company_name(text)
   
    if PLACEHOLDER_TOKEN in protected_text and protected_text.strip() == PLACEHOLDER_TOKEN:
        final_text = CORRECT_COMPANY_NAME
        translated_texts.append(final_text)
        translation_log.append("Company name only - no translation")
        print(f"{original_text} -> {final_text}")
        return final_text
    # Company name only (no context), restore without translating
    elif PLACEHOLDER_TOKEN in protected_text:
        translated = translate_chunk(protected_text, source_lang, target_lang)
        final_text = restore_company_name(translated)
        translated_texts.append(final_text)
        translation_log.append("Translated (with placeholder restored)")
        print(f"{original_text} -> {final_text}")
        return final_text
 
 
    # if protected_text.strip() == PLACEHOLDER_TOKEN:
    #     final_text = CORRECT_COMPANY_NAME
    #     translated_texts.append(final_text)
    #     translation_log.append("Company name only - no translation")
    #     print(f"{original_text} -> {final_text}")
    #     return final_text
 
    chunks = split_text_into_chunks(protected_text)
    translated_chunks = []
 
    with ThreadPoolExecutor(max_workers=5) as executor:
        future_to_chunk = {executor.submit(translate_chunk, chunk, source_lang, target_lang): chunk for chunk in chunks if chunk.strip()}
        for future in as_completed(future_to_chunk):
            try:
                result = future.result()
                if result is None:
                    result = future_to_chunk[future]
                translated_chunks.append(result)
            except Exception as e:
                print(f"Translation error: {e}", file=sys.stderr)
                result = future_to_chunk[future]
                translated_chunks.append(result)
 
    translated_text = ' '.join(translated_chunks)
    final_text = restore_company_name(translated_text)
 
    translated_texts.append(final_text)
    translation_log.append("Translated (company name preserved)" if final_text.strip() != original_text.strip() else "No change after translation")
    print(f"{original_text} -> {final_text}")
    return final_text
 
def process_text_frame(text_frame, source_lang, target_lang):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            if original_text.strip():
                translated_text = translate_text_safely(original_text, source_lang, target_lang)
                run.text = translated_text
 
# Utility to estimate if text fits in a shape (approximate, since pptx has no direct text measurement)
def text_fits_in_shape(text_frame, shape_width, shape_height, font_size_pt):
    # Estimate: average character width is 0.5 * font size, height is font size
    # This is a rough approximation for Latin/Cyrillic scripts
    lines = text_frame.text.split('\n')
    max_line_len = max((len(line) for line in lines), default=0)
    est_text_width = max_line_len * font_size_pt * 0.5
    est_text_height = len(lines) * font_size_pt * 1.2
    return est_text_width <= shape_width and est_text_height <= shape_height
 
# Adjust font size in a text frame to fit the shape
 
def fit_text_frame_font(text_frame, shape_width, shape_height):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                font_size = run.font.size.pt
            else:
                font_size = 18  # default
            # Reduce font size until text fits or minimum size reached
            while font_size > MIN_FONT_SIZE_PT:
                run.font.size = Pt(font_size)
                if text_fits_in_shape(text_frame, shape_width, shape_height, font_size):
                    break
                font_size -= 1
            run.font.size = Pt(max(font_size, MIN_FONT_SIZE_PT))
 
# Update process_shape_text to fit text after translation
 
def process_shape_text(shape, source_lang, target_lang):
    # Process any shape with a text_frame, even if text is empty
    if hasattr(shape, 'text_frame') and shape.text_frame is not None:
        process_text_frame(shape.text_frame, source_lang, target_lang)
        # Fit text after translation
        fit_text_frame_font(shape.text_frame, shape.width, shape.height)
 
    # Process text in table cells if it's a table
    if hasattr(shape, 'has_table') and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if hasattr(cell, 'text_frame') and cell.text_frame is not None:
                    process_text_frame(cell.text_frame, source_lang, target_lang)
                    fit_text_frame_font(cell.text_frame, cell.width, cell.height)
 
    # Recursively process nested shapes (e.g., grouped shapes, SmartArt)
    if hasattr(shape, 'shapes') and shape.shapes:
        for nested_shape in shape.shapes:
            process_shape_text(nested_shape, source_lang, target_lang)
 
 
def scale_slide_fonts(slide, scale=0.8, default_size_pt=18):
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame is not None:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    current_size = run.font.size.pt if run.font.size else default_size_pt
                    new_size = max(current_size * scale, MIN_FONT_SIZE_PT)
                    run.font.size = Pt(new_size)
        if hasattr(shape, 'has_table') and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, 'text_frame') and cell.text_frame is not None:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                current_size = run.font.size.pt if run.font.size else default_size_pt
                                new_size = max(current_size * scale, MIN_FONT_SIZE_PT)
                                run.font.size = Pt(new_size)
 
def translate_presentation(input_file, output_file, target_lang):
    source_lang = 'auto'
    prs = Presentation(input_file)
    scaled_slides = []
    image_slides = []
 
    for i, slide in enumerate(prs.slides):
        should_scale = False
        for shape in slide.shapes:
            # if hasattr(shape, 'text') and shape.text.strip():
            process_shape_text(shape, source_lang, target_lang)
            try:
                text_length = len(shape.text)
                shape_width_in_inches = shape.width.inches if hasattr(shape, 'width') else 10.0
                if text_length > 200 and shape_width_in_inches < 5.0:
                    print(f"Slide {i+1}: Long text ({text_length} chars) in narrow width ({shape_width_in_inches:.2f} inches)")
                    should_scale = True
            except Exception as e:
                print(f"Error evaluating scaling condition: {e}")
 
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if (i + 1) not in image_slides:
                    image_slides.append(i + 1)
 
        if should_scale:
            print(f"Scaling down fonts for slide {i+1}...")
            scale_slide_fonts(slide)
            scaled_slides.append(i + 1)
 
    prs.save(output_file)
    print("Translation complete. Saved to:", output_file)
 
    base_output = os.path.splitext(os.path.basename(output_file))[0]
    log_filename = f"translation_log_{base_output}.txt"
    log_path = os.path.join(os.path.dirname(output_file), log_filename)
    try:
        with open(log_path, "w", encoding="utf-8") as f_log:
            f_log.write("Translation Summary Log\n")
            f_log.write("=======================\n")
            f_log.write(f"Output File: {output_file}\n")
            f_log.write(f"Time: {datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}\n\n")
            if scaled_slides:
                f_log.write(f"Slides with font scaling applied: {scaled_slides}\n")
            else:
                f_log.write("No slides required font scaling.\n")
            if image_slides:
                f_log.write(f"Slides containing images: {image_slides}\n")
            else:
                f_log.write("No image-containing slides detected.\n")
        print(f"\nSummary saved to {log_path}")
    except Exception as e:
        print(f"Failed to write summary log: {e}")
 
def main():
    if len(sys.argv) != 4:
        print(f"Usage: python {sys.argv[0]} <input_file> <output_file> <target_lang>")
        sys.exit(1)
    translate_presentation(sys.argv[1], sys.argv[2], sys.argv[3])
 
if __name__ == "__main__":
    main()
