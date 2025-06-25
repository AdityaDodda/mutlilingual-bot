# Source lang - target lang
from deep_translator import GoogleTranslator
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, as_completed
import sys

# --- Helper functions (no changes needed here) ---

def split_text_into_chunks(text, max_length=2000):
    """Split text into chunks to handle the translator's character limit."""
    return [text[i:i+max_length] for i in range(0, len(text), max_length)]

def translate_chunk(chunk, source_lang, target_lang):
    # Use 'auto' for source language to let the library detect it
    return GoogleTranslator(source=source_lang, target=target_lang).translate(chunk)

def translate_text_safely(text, source_lang, target_lang):
    """Safely translate text, handling errors and empty strings."""
    if not text or not text.strip():
        return text

    chunks = split_text_into_chunks(text)
    translated_chunks = []

    with ThreadPoolExecutor(max_workers=5) as executor:
        future_to_chunk = {executor.submit(translate_chunk, chunk, source_lang, target_lang): chunk for chunk in chunks if chunk.strip()}
        for future in as_completed(future_to_chunk):
            try:
                result = future.result()
                if result is None:
                    # If translation fails and returns None, use original chunk
                    result = future_to_chunk[future]
                translated_chunks.append(result)
            except Exception as e:
                print(f"Translation error: {e}", file=sys.stderr)
                # On error, fall back to the original text chunk
                translated_chunks.append(future_to_chunk[future])

    return ' '.join(translated_chunks)


# --- PPTX Processing functions (no changes needed here) ---

def process_shape_text(shape, source_lang, target_lang, depth=0):
    indent = "  " * depth
    shape_name = getattr(shape, 'name', f"Shape{id(shape)}")

    if hasattr(shape, 'text_frame') and shape.text.strip():
        # print(f"{indent}Processing text in shape: {shape_name}")
        process_text_frame(shape.text_frame, source_lang, target_lang)

    if hasattr(shape, 'has_table') and shape.has_table:
        # print(f"{indent}Processing table in shape: {shape_name}")
        for row in shape.table.rows:
            for cell in row.cells:
                if hasattr(cell, 'text_frame') and cell.text.strip():
                    process_text_frame(cell.text_frame, source_lang, target_lang)

    if hasattr(shape, 'shapes'):
        # print(f"{indent}Found {len(shape.shapes)} nested shapes in: {shape_name}")
        for nested_shape in shape.shapes:
            process_shape_text(nested_shape, source_lang, target_lang, depth + 1)

def process_text_frame(text_frame, source_lang, target_lang):
    for paragraph in text_frame.paragraphs:
        if paragraph.text.strip():
            original_runs = list(paragraph.runs)
            original_text = paragraph.text
            translated_text = translate_text_safely(original_text, source_lang, target_lang)

            # Clear existing runs
            while len(paragraph.runs) > 0:
                paragraph._p.remove(paragraph.runs[0]._r)

            # Add a new run with the translated text
            new_run = paragraph.add_run()
            new_run.text = translated_text

            # Attempt to copy formatting from the first original run
            if original_runs:
                first_run = original_runs[0]
                try:
                    if hasattr(first_run, 'font'):
                        if hasattr(first_run.font, 'name') and first_run.font.name:
                            new_run.font.name = first_run.font.name
                        if hasattr(first_run.font, 'size') and first_run.font.size is not None:
                            new_run.font.size = first_run.font.size
                        if hasattr(first_run.font, 'bold'):
                            new_run.font.bold = first_run.font.bold
                        if hasattr(first_run.font, 'italic'):
                            new_run.font.italic = first_run.font.italic
                        if hasattr(first_run.font, 'underline'):
                            new_run.font.underline = first_run.font.underline
                        if hasattr(first_run.font.color, 'rgb') and first_run.font.color.rgb is not None:
                            new_run.font.color.rgb = first_run.font.color.rgb
                        elif hasattr(first_run.font.color, 'theme_color') and first_run.font.color.theme_color is not None:
                            new_run.font.color.theme_color = first_run.font.color.theme_color
                except Exception as e:
                    print(f"Warning: Could not copy all formatting properties: {e}", file=sys.stderr)

def translate_presentation(input_file, output_file, target_lang):
    # MODIFIED: Use 'auto' to automatically detect the source language.
    # This is much more flexible than hardcoding a language.
    source_lang = 'auto'
    
    print(f"Loading presentation: {input_file}")
    print(f"Using auto-detected source language to translate to target: '{target_lang}'")
    
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    print(f"Found {total_slides} slides to translate.")

    for i, slide in enumerate(prs.slides):
        print(f"Translating slide {i+1}/{total_slides}...")
        for shape_idx, shape in enumerate(slide.shapes):
            # The smaller logs are commented out to keep the output clean
            # print(f"  Processing shape {shape_idx+1}/{len(slide.shapes)}")
            process_shape_text(shape, source_lang, target_lang)
            
    print(f"Saving translated presentation to: {output_file}")
    prs.save(output_file)
    print("Translation complete!")


def main():
    print("--- Python Conversion Script Started ---")
    print(f"Raw command-line arguments: {sys.argv}")
    
    # MODIFIED: Check for the correct number of arguments
    if len(sys.argv) != 4:
        # MODIFIED: Use sys.argv[0] to show the correct script name in the usage message
        print(f"Error: Invalid arguments.", file=sys.stderr)
        print(f"Usage: python {sys.argv[0]} <input_file> <output_file> <target_lang>", file=sys.stderr)
        sys.exit(1)
        
    # MODIFIED: Get arguments from the command line
    pptx_file_path = sys.argv[1]
    output_file_path = sys.argv[2]
    target_lang = sys.argv[3]

    print(f"Input file: {pptx_file_path}")
    print(f"Output file: {output_file_path}")
    print(f"Target language: {target_lang}")
    print("------------------------------------")
    
    try:
        translate_presentation(
            input_file=pptx_file_path,
            output_file=output_file_path,
            target_lang=target_lang
        )
    except Exception as e:
        print(f"An unexpected error occurred: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()